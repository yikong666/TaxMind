"""多格式文档文本与标题结构提取。"""

# 各格式最终归一为标题和正文，图片内容通过 OCR 提取。
import re
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import fitz
import numpy as np
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation


class DocumentParseError(ValueError):
    pass


@dataclass(frozen=True)
class ParsedSection:
    heading: str | None
    content: str


class DocumentParser:
    def parse(self, filename: str, content: bytes) -> list[ParsedSection]:
        extension = Path(filename).suffix.lower()
        parser = {
            ".md": self._parse_markdown,
            ".txt": self._parse_text,
            ".html": self._parse_html,
            ".htm": self._parse_html,
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".pptx": self._parse_pptx,
            ".ppt": self._unsupported_legacy_office,
            ".doc": self._unsupported_legacy_office,
            ".jpg": self._parse_image,
            ".jpeg": self._parse_image,
            ".png": self._parse_image,
            ".webp": self._parse_image,
        }.get(extension)
        if parser is None:
            raise DocumentParseError("不支持的文档格式")
        sections = parser(content)
        cleaned = [
            ParsedSection(item.heading, self._clean_text(item.content))
            for item in sections
            if self._clean_text(item.content)
        ]
        if not cleaned:
            raise DocumentParseError("文档中未提取到有效文本")
        return cleaned

    def _parse_markdown(self, content: bytes) -> list[ParsedSection]:
        text = self._decode(content)
        sections: list[ParsedSection] = []
        heading: str | None = None
        buffer: list[str] = []
        for line in text.splitlines():
            match = re.match(r"^#{1,6}\s+(.+?)\s*$", line)
            if match:
                if buffer:
                    sections.append(ParsedSection(heading, "\n".join(buffer)))
                heading, buffer = match.group(1), []
            else:
                buffer.append(line)
        if buffer:
            sections.append(ParsedSection(heading, "\n".join(buffer)))
        return sections

    def _parse_text(self, content: bytes) -> list[ParsedSection]:
        return [ParsedSection(None, self._decode(content))]

    def _parse_html(self, content: bytes) -> list[ParsedSection]:
        soup = BeautifulSoup(content, "html.parser")
        for element in soup(["script", "style", "noscript"]):
            element.decompose()
        sections: list[ParsedSection] = []
        heading: str | None = None
        buffer: list[str] = []
        for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]):
            text = element.get_text(" ", strip=True)
            if not text:
                continue
            if element.name.startswith("h"):
                if buffer:
                    sections.append(ParsedSection(heading, "\n".join(buffer)))
                heading, buffer = text, []
            else:
                buffer.append(text)
        if buffer:
            sections.append(ParsedSection(heading, "\n".join(buffer)))
        return sections

    def _parse_pdf(self, content: bytes) -> list[ParsedSection]:
        with fitz.open(stream=content, filetype="pdf") as pdf:
            return [
                ParsedSection(f"第 {index + 1} 页", page.get_text("text"))
                for index, page in enumerate(pdf)
            ]

    def _parse_docx(self, content: bytes) -> list[ParsedSection]:
        document = DocxDocument(BytesIO(content))
        sections: list[ParsedSection] = []
        heading: str | None = None
        buffer: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            if paragraph.style and paragraph.style.name.lower().startswith("heading"):
                if buffer:
                    sections.append(ParsedSection(heading, "\n".join(buffer)))
                heading, buffer = text, []
            else:
                buffer.append(text)
        if buffer:
            sections.append(ParsedSection(heading, "\n".join(buffer)))
        return sections

    def _parse_pptx(self, content: bytes) -> list[ParsedSection]:
        presentation = Presentation(BytesIO(content))
        sections: list[ParsedSection] = []
        for index, slide in enumerate(presentation.slides):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                sections.append(
                    ParsedSection(
                        texts[0] or f"第 {index + 1} 页",
                        "\n".join(texts[1:] or texts),
                    )
                )
        return sections

    def _parse_image(self, content: bytes) -> list[ParsedSection]:
        from rapidocr_onnxruntime import RapidOCR

        image = np.asarray(Image.open(BytesIO(content)).convert("RGB"))
        result, _ = RapidOCR()(image)
        text = "\n".join(item[1] for item in (result or []) if len(item) > 1)
        return [ParsedSection("图片 OCR", text)]

    @staticmethod
    def _unsupported_legacy_office(_: bytes) -> list[ParsedSection]:
        raise DocumentParseError("旧版 DOC/PPT 暂不支持解析，请转换为 DOCX/PPTX 后重新上传")

    @staticmethod
    def _decode(content: bytes) -> str:
        for encoding in ("utf-8-sig", "utf-8", "gb18030"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise DocumentParseError("文本编码无法识别，请转换为 UTF-8")

    @staticmethod
    def _clean_text(text: str) -> str:
        text = text.replace("\x00", "")
        lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
        return "\n".join(line for line in lines if line).strip()
